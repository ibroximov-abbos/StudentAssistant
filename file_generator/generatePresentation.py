from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

import json
import re
def extract_json_regex(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    pattern = r'```json\s*(.*?)\s*```'
    match = re.search(pattern, content, re.DOTALL)
    if match:
        json_str = match.group(1)
        data = json.loads(json_str)
        return data
    return None

def create_presentation(theme, full_name, username):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    background = slide1.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    fill.gradient_stops[0].color.rgb = RGBColor(0xE0, 0xF7, 0xFA)
    fill.gradient_stops[1].color.rgb = RGBColor(0xB2, 0xEB, 0xF2)
    height = Inches(4)
    width = Inches(0.08)

    border_line = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    left=Inches(0.8), top=Inches(2),  # x, y pozitsiyasi
    width=width, height=height  # kenglik, balandlik
    )
    border_line.fill.solid()
    border_line.fill.fore_color.rgb = RGBColor(0x00, 0x83, 0x8F)
    border_line.line.fill.background()

    title_box = slide1.shapes.add_textbox(
    Inches(1.2), Inches(2),  # x, y
    Inches(8), Inches(2)     # kenglik, balandlik
    )
    title_frame = title_box.text_frame
    title_frame.text = theme
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36) 
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x00, 0x60, 0x64)
    title_para.alignment = PP_ALIGN.LEFT
    title_para.line_spacing = 1.1

    author_box = slide1.shapes.add_textbox(
    Inches(1.2), Inches(5.8),  # x, y
    Inches(8), Inches(0.5)      # kenglik, balandlik
    )
    author_frame = author_box.text_frame
    author_frame.text = f"Taqdimotchi: {full_name}"

    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(17)  # 1.2em ≈ 17pt
    author_para.font.color.rgb = RGBColor(0x45, 0x5A, 0x64)
    author_para.alignment = PP_ALIGN.LEFT

    data = extract_json_regex(f"{username}.txt")
    
    if data:
        for key in data:
            print(key)
            blank_slide_layout = prs.slide_layouts[6]  # Bo'sh layout
            slide = prs.slides.add_slide(blank_slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xED, 0xE7, 0xF6)
            
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), 
                Inches(9), Inches(0.6)
            )
            title_frame = title_box.text_frame
            title_frame.text = data[key]['title']

            len_title = len(data[key]['title'])
            if  len_title <= 45:
                title_frame.paragraphs[0].font.size = Pt(32)
            elif 45 < len_title <= 60:
                title_frame.paragraphs[0].font.size = Pt(28)
            else:
                title_frame.paragraphs[0].font.size = Pt(24)

            p = title_frame.paragraphs[0]
            p.font.color.rgb = RGBColor(0x0D, 0x13, 0x47)  
            p.font.bold = True
            pointer1 = data[key]['content']['pointer1']
            pointer2 = data[key]['content']['pointer2']
            pointer3 = data[key]['content']['pointer3']

            if int(key) % 3 == 0:
                pointer2 += f"\n{pointer3}"
            elif int(key) % 3 == 2:
                pointer1 += f"\n\n{pointer2}\n{pointer3}"


            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(7)
            shape = slide.shapes.add_shape(
             MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
            )
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x28, 0x35, 0x93)
            line = shape.line
            line.color.rgb = RGBColor(0x0D, 0x13, 0x47)
            line.width = Pt(2)

            text_frame = shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            text_frame.margin_top = Inches(0.3)
            text_frame.margin_left = Inches(0.3)
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            p1 = text_frame.paragraphs[0]
            run1 = p1.add_run()
            run1.text = pointer1
            run1.font.color.rgb = RGBColor(0xED, 0xE7, 0xF6)
            font1 = run1.font
            font1.name = 'Arial' 
            font1.size = Pt(18)
            font1.bold = True

            if int(key) % 3 == 0:
                left = Inches(0.5)
                top = Inches(4)
                width = Inches(9)
                height = Inches(3)
                shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
                )
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0x28, 0x35, 0x93) 
                line = shape.line
                line.color.rgb = RGBColor(0x0D, 0x13, 0x47)
                line.width = Pt(1.5) 

                text_frame = shape.text_frame
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
                text_frame.margin_top = Inches(0.3)
                text_frame.margin_left = Inches(0.3)
                text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

                p1 = text_frame.paragraphs[0]
                run1 = p1.add_run()
                run1.text = pointer2
                font1 = run1.font
                font1.name = 'Arial' # Yoki boshqa shrift
                font1.size = Pt(18)
                font1.bold = True
                font1.color.rgb = RGBColor(0xED, 0xE7, 0xF6)

            if int(key) % 3 == 1:
                left = Inches(0.5)
                top = Inches(4)
                width = Inches(4)
                height = Inches(3)
                shape1 = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
                )
                shape2 = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), top, width, height
                )
                for shape in (shape1, shape2):
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(0x28, 0x35, 0x93)
                    line = shape.line
                    line.color.rgb = RGBColor(0x0D, 0x13, 0x47)
                    line.width = Pt(1.5) # Chiziq qalinligi

                    text_frame = shape.text_frame
                    text_frame.vertical_anchor = MSO_ANCHOR.TOP
                    text_frame.margin_top = Inches(0.3)
                    text_frame.margin_left = Inches(0.3)
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

                    p1 = text_frame.paragraphs[0]
                    run1 = p1.add_run()
                    if shape == shape1:
                        run1.text = pointer2
                    else:
                        run1.text = pointer3
                    font1 = run1.font
                    font1.name = 'Arial' # Yoki boshqa shrift
                    font1.size = Pt(18)
                    font1.bold = True
                    font1.color.rgb = RGBColor(0xED, 0xE7, 0xF6)

    prs.save(f'{username}.pptx')
    print(f"{username}.pptx yaratildi")

